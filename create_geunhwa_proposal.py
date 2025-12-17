#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
경주 근화여자중학교 맞춤형 겨울방학 자기주도학습 프로그램 제안서 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle=""):
    """타이틀 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경색 (여학교 이미지 - 우아한 보라/핑크)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(123, 31, 162)  # 보라색
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 부제목
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(26)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items, emoji=""):
    """내용 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 제목 영역
    title_shape = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(123, 31, 162)
    title_shape.line.color.rgb = RGBColor(123, 31, 162)
    
    # 제목 텍스트
    title_frame = title_shape.text_frame
    title_frame.text = f"{emoji} {title}" if emoji else title
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 내용 영역
    content_top = 1.8
    for i, item in enumerate(content_items):
        text_box = slide.shapes.add_textbox(
            Inches(0.8), 
            Inches(content_top + i * 0.7), 
            Inches(8.4), 
            Inches(0.6)
        )
        text_frame = text_box.text_frame
        text_frame.text = item
        para = text_frame.paragraphs[0]
        para.font.size = Pt(24)
        para.font.color.rgb = RGBColor(51, 51, 51)
        para.space_before = Pt(6)
        para.space_after = Pt(6)

def add_program_detail_slide(prs, title, programs, emoji=""):
    """프로그램 상세 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # 제목
    title_shape = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(123, 31, 162)
    title_shape.line.color.rgb = RGBColor(123, 31, 162)
    
    title_frame = title_shape.text_frame
    title_frame.text = f"{emoji} {title}" if emoji else title
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(38)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 프로그램 카드들
    card_top = 1.8
    for program in programs:
        # 카드 배경
        card = slide.shapes.add_shape(
            1,
            Inches(0.8), Inches(card_top), Inches(8.4), Inches(1.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(200, 200, 200)
        card.line.width = Pt(1)
        
        # 프로그램 제목
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(card_top + 0.15), Inches(8), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = program['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(26)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(123, 31, 162)
        
        # 프로그램 설명
        desc_box = slide.shapes.add_textbox(
            Inches(1), Inches(card_top + 0.6), Inches(8), Inches(0.6)
        )
        desc_frame = desc_box.text_frame
        desc_frame.text = program['desc']
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(20)
        desc_para.font.color.rgb = RGBColor(80, 80, 80)
        
        card_top += 1.5

def add_contact_slide(prs):
    """연락처 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(243, 229, 245)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "📞 문의 및 상담"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(123, 31, 162)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 연락처 정보
    contact_items = [
        "📱 전화: 010-2657-3481",
        "👤 담당자: 정라미",
        "🏫 제안 대상: 경주 근화여자중학교",
        "📚 프로그램: 1,2학년 겨울방학 자기주도학습",
        "",
        "💡 학교 맞춤형 프로그램 운영",
        "💡 체계적인 학습관리 시스템"
    ]
    
    contact_top = 3
    for i, item in enumerate(contact_items):
        text_box = slide.shapes.add_textbox(
            Inches(2), 
            Inches(contact_top + i * 0.45), 
            Inches(6), 
            Inches(0.4)
        )
        text_frame = text_box.text_frame
        text_frame.text = item
        para = text_frame.paragraphs[0]
        para.font.size = Pt(22) if item else Pt(12)
        para.font.color.rgb = RGBColor(51, 51, 51)
        para.alignment = PP_ALIGN.CENTER

def create_presentation():
    """프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. 표지
    add_title_slide(
        prs, 
        "경주 근화여자중학교",
        "1,2학년 겨울방학 자기주도학습 프로그램 제안서"
    )
    
    # 2. 제안 개요
    add_content_slide(
        prs,
        "제안 개요",
        [
            "🏫 대상: 경주 근화여자중학교 1학년, 2학년",
            "❄️ 기간: 겨울방학 30일 집중 프로그램",
            "🎯 목표: 자기주도학습 능력 향상 및 학업 성취도 제고",
            "👩‍🏫 방식: 학교 맞춤형 체계적 학습관리",
            "📊 특징: 개인별 맞춤 지도 및 진도 관리",
            "💻 시스템: 온라인 학습관리 플랫폼 제공"
        ],
        "📋"
    )
    
    # 3. 프로그램의 필요성
    add_content_slide(
        prs,
        "프로그램의 필요성",
        [
            "✓ 겨울방학 학습 공백 방지 및 학력 향상",
            "✓ 자기주도학습 습관 형성의 중요한 시기",
            "✓ 고등학교 진학 준비 및 학업 기초 다지기",
            "✓ 체계적인 학습관리를 통한 성적 향상",
            "✓ 개인별 맞춤 교육으로 학습 효율 극대화",
            "✓ 진로 탐색 및 목표 설정 기회 제공"
        ],
        "💡"
    )
    
    # 4. 프로그램 특징
    add_content_slide(
        prs,
        "근화여중 맞춤형 프로그램 특징",
        [
            "🎯 여학생 특성을 고려한 맞춤형 학습 지도",
            "📚 30일 집중 학습으로 학습 습관 완성",
            "👩‍🏫 1:1 개인별 학습관리 및 피드백",
            "📊 체계적인 진도 관리 및 성취도 평가",
            "💻 온라인 학습관리 시스템 활용",
            "📞 학부모 상담 및 정기 학습 리포트 제공",
            "🎓 고등학교 준비 및 진로 지도"
        ],
        "✨"
    )
    
    # 5. 1학년 프로그램
    add_program_detail_slide(
        prs,
        "1학년 프로그램",
        [
            {
                'title': '📖 교과 학습 강화',
                'desc': '국어, 영어, 수학, 사회, 과학 등 주요 교과 집중 학습 및 내신 대비'
            },
            {
                'title': '✍️ 자기주도학습 훈련',
                'desc': '학습 계획 수립, 시간 관리, 효과적인 공부 방법 습득'
            },
            {
                'title': '📝 기초 학력 완성',
                'desc': '취약 과목 보완 및 학습 기초 다지기'
            }
        ],
        "📚"
    )
    
    # 6. 2학년 프로그램
    add_program_detail_slide(
        prs,
        "2학년 프로그램",
        [
            {
                'title': '📖 교과 심화 학습',
                'desc': '주요 교과 심화 학습 및 고등학교 준비 선행 학습'
            },
            {
                'title': '🎯 고등학교 진학 준비',
                'desc': '진로 탐색, 목표 설정, 고등 과정 미리보기'
            },
            {
                'title': '💪 학습 역량 강화',
                'desc': '심화 문제 해결 능력 향상 및 자기주도학습 완성'
            }
        ],
        "📚"
    )
    
    # 7. 학습관리 시스템
    add_content_slide(
        prs,
        "체계적인 학습관리 시스템",
        [
            "📅 일일 학습 목표 설정 및 진도 체크",
            "✍️ 학습 일지 작성 및 자기 점검",
            "📊 주간 학습 평가 및 성취도 분석",
            "👩‍🏫 개인별 학습 상담 및 피드백",
            "📞 학부모 주간 리포트 및 정기 상담",
            "💻 온라인 플랫폼을 통한 학습 관리",
            "📈 모의고사 및 실력 진단 평가"
        ],
        "⚙️"
    )
    
    # 8. 프로그램 운영 계획
    add_content_slide(
        prs,
        "프로그램 운영 계획",
        [
            "📅 운영 기간: 겨울방학 중 30일 (협의 후 확정)",
            "⏰ 운영 시간: 평일 09:00 ~ 17:00 (점심시간 포함)",
            "👥 운영 인원: 학년별 소그룹 또는 개인별 맞춤",
            "📍 운영 장소: 학교 또는 DA.UM 학습센터",
            "📚 교재 및 자료: 학교 교과서 + 맞춤형 학습 자료",
            "👩‍🏫 지도 교사: 과목별 전문 강사진",
            "💻 온라인 지원: 학습관리 앱 및 온라인 상담"
        ],
        "📋"
    )
    
    # 9. 기대 효과
    add_content_slide(
        prs,
        "기대 효과",
        [
            "🌟 자기주도학습 습관 완성",
            "📈 학업 성취도 및 내신 성적 향상",
            "💪 학습 자신감 및 동기 부여",
            "🎯 명확한 학습 목표 및 진로 방향 설정",
            "📚 고등학교 진학 준비 완료",
            "👩‍👧 학부모와의 긴밀한 소통 및 협력",
            "✨ 근화여중 학생들의 학업 경쟁력 강화"
        ],
        "✨"
    )
    
    # 10. 연락처
    add_contact_slide(prs)
    
    # 11. 감사 슬라이드
    add_title_slide(
        prs,
        "감사합니다",
        "근화여자중학교 학생들의 성장을 함께 만들어가겠습니다"
    )
    
    # 저장
    output_file = '/home/user/webapp/근화여자중학교_겨울방학_자기주도학습_프로그램_제안서.pptx'
    prs.save(output_file)
    print(f"✅ 제안서 생성 완료: {output_file}")
    return output_file

if __name__ == "__main__":
    create_presentation()
