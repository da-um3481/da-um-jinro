#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DA.UM 겨울방학 프로그램 학교 설명회 자료 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title, subtitle=""):
    """타이틀 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    
    # 배경색 설정 (진한 파란색)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(41, 98, 255)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 부제목
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items, emoji=""):
    """내용 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경색 (흰색)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 제목 영역 (파란색 배경)
    title_shape = slide.shapes.add_shape(
        1,  # 사각형
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(41, 98, 255)
    title_shape.line.color.rgb = RGBColor(41, 98, 255)
    
    # 제목 텍스트
    title_frame = title_shape.text_frame
    title_frame.text = f"{emoji} {title}" if emoji else title
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 내용 영역
    content_top = 1.8
    for i, item in enumerate(content_items):
        text_box = slide.shapes.add_textbox(
            Inches(0.8), 
            Inches(content_top + i * 0.7), 
            Inches(8.4), 
            Inches(0.6)
        )
        text_frame = text_box.text_frame
        text_frame.text = item
        para = text_frame.paragraphs[0]
        para.font.size = Pt(24)
        para.font.color.rgb = RGBColor(51, 51, 51)
        para.space_before = Pt(6)
        para.space_after = Pt(6)

def add_program_detail_slide(prs, title, programs, emoji=""):
    """프로그램 상세 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # 제목
    title_shape = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(41, 98, 255)
    title_shape.line.color.rgb = RGBColor(41, 98, 255)
    
    title_frame = title_shape.text_frame
    title_frame.text = f"{emoji} {title}" if emoji else title
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(38)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 프로그램 카드들
    card_top = 1.8
    for program in programs:
        # 카드 배경
        card = slide.shapes.add_shape(
            1,
            Inches(0.8), Inches(card_top), Inches(8.4), Inches(1.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(200, 200, 200)
        card.line.width = Pt(1)
        
        # 프로그램 제목
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(card_top + 0.15), Inches(8), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = program['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(26)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(41, 98, 255)
        
        # 프로그램 설명
        desc_box = slide.shapes.add_textbox(
            Inches(1), Inches(card_top + 0.6), Inches(8), Inches(0.6)
        )
        desc_frame = desc_box.text_frame
        desc_frame.text = program['desc']
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(20)
        desc_para.font.color.rgb = RGBColor(80, 80, 80)
        
        card_top += 1.5

def add_contact_slide(prs):
    """연락처 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경색 (연한 파란색)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 240, 255)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "📞 문의 및 상담"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(41, 98, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 연락처 정보
    contact_items = [
        "📱 전화: 010-2657-3481",
        "👤 담당자: 정라미",
        "🌐 웹사이트: da-um3481.github.io/da-um-jinro",
        "",
        "💡 학생 맞춤형 진로진학 컨설팅",
        "💡 체계적인 학습관리 시스템"
    ]
    
    contact_top = 3
    for i, item in enumerate(contact_items):
        text_box = slide.shapes.add_textbox(
            Inches(2), 
            Inches(contact_top + i * 0.5), 
            Inches(6), 
            Inches(0.4)
        )
        text_frame = text_box.text_frame
        text_frame.text = item
        para = text_frame.paragraphs[0]
        para.font.size = Pt(24) if item else Pt(12)
        para.font.color.rgb = RGBColor(51, 51, 51)
        para.alignment = PP_ALIGN.CENTER

def create_presentation():
    """프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. 표지
    add_title_slide(
        prs, 
        "DA.UM 다움진로진학컨설팅",
        "겨울방학 프로그램 안내"
    )
    
    # 2. 소개 슬라이드
    add_content_slide(
        prs,
        "DA.UM 소개",
        [
            "✓ 전문 진로진학 컨설팅 기관",
            "✓ 학생 개인별 맞춤 학습관리",
            "✓ 고등학생 및 중학생 전문 프로그램",
            "✓ 체계적인 30일/100일 학습 시스템",
            "✓ 겨울방학 집중 프로그램 운영"
        ],
        "🎯"
    )
    
    # 3. 겨울방학 프로그램 특징
    add_content_slide(
        prs,
        "겨울방학 프로그램 특징",
        [
            "❄️ 30일 집중 학습 프로그램",
            "📊 학년별 맞춤형 커리큘럼",
            "👨‍🏫 1:1 개인별 학습관리",
            "📈 체계적인 진도 관리 및 피드백",
            "🎓 진로진학 맞춤 컨설팅",
            "💻 온라인 학습관리 시스템 제공"
        ],
        "❄️"
    )
    
    # 4. 고등학생 프로그램
    add_program_detail_slide(
        prs,
        "고등학생 프로그램",
        [
            {
                'title': '📚 고1 프로그램',
                'desc': '기초 학력 다지기 + 내신 대비 + 학습 습관 형성'
            },
            {
                'title': '📚 고2 프로그램',
                'desc': '심화 학습 + 수능 기초 준비 + 진로 탐색'
            },
            {
                'title': '📚 고3 프로그램',
                'desc': '수능 집중 대비 + 입시 전략 수립 + 최종 마무리'
            }
        ],
        "🎓"
    )
    
    # 5. 중학생 프로그램
    add_program_detail_slide(
        prs,
        "중학생 이하 프로그램",
        [
            {
                'title': '📖 중학교 1~3학년',
                'desc': '교과 학습 관리 + 내신 대비 + 학습 역량 강화'
            },
            {
                'title': '📖 예비 중학생 (초6)',
                'desc': '중학교 준비 + 선행 학습 + 학습 습관 형성'
            }
        ],
        "📖"
    )
    
    # 6. 학습관리 시스템
    add_content_slide(
        prs,
        "학습관리 시스템",
        [
            "⚙️ 일일 학습 진도 체크",
            "⚙️ 주간 학습 리포트 제공",
            "⚙️ 온라인 학습 플랫폼 활용",
            "⚙️ 학부모 상담 및 피드백",
            "⚙️ 개인별 취약점 분석 및 보완",
            "⚙️ 모의고사 및 평가 관리"
        ],
        "⚙️"
    )
    
    # 7. 프로그램 구성
    add_content_slide(
        prs,
        "프로그램 구성",
        [
            "📅 겨울방학 30일 집중 프로그램",
            "📅 학기중 100일 정규 프로그램",
            "🎯 단계별 학습 목표 설정",
            "🎯 정기적인 학습 점검 및 평가",
            "🎯 맞춤형 학습 자료 제공",
            "🎯 진로진학 개별 컨설팅"
        ],
        "📋"
    )
    
    # 8. 기대 효과
    add_content_slide(
        prs,
        "프로그램 기대 효과",
        [
            "🌟 체계적인 학습 습관 형성",
            "🌟 학업 성취도 향상",
            "🌟 자기주도 학습 능력 배양",
            "🌟 명확한 진로 목표 설정",
            "🌟 입시 경쟁력 강화",
            "🌟 학습 동기 부여 및 자신감 향상"
        ],
        "✨"
    )
    
    # 9. 연락처
    add_contact_slide(prs)
    
    # 10. 감사 슬라이드
    add_title_slide(
        prs,
        "감사합니다",
        "DA.UM과 함께 성공적인 겨울방학을 만들어가세요"
    )
    
    # 저장
    output_file = '/home/user/webapp/DA_UM_겨울방학_프로그램_설명회.pptx'
    prs.save(output_file)
    print(f"✅ 프레젠테이션 생성 완료: {output_file}")
    return output_file

if __name__ == "__main__":
    create_presentation()
